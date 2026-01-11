#!/usr/bin/env python3
"""Import a PPTX as a course: converts slides to PNG, extracts notes & assets, and inserts into SQLite DB."""
import argparse
from pathlib import Path
import subprocess
import json
import shutil
import sys

from sqlalchemy import create_engine
from sqlalchemy.orm import sessionmaker
from models import Base, Course, Lesson

try:
    from pptx import Presentation
except Exception:
    Presentation = None


def convert_pptx(pptx_path: Path, out_dir: Path):
    # call existing convert.py script in project if present
    convert_script = Path(__file__).parent.parent / 'pptx_viewer' / 'convert.py'
    if convert_script.exists():
        cmd = [sys.executable, str(convert_script), '--pptx', str(pptx_path), '--out', str(out_dir), '--mode', 'pdf']
        print('Running:', ' '.join(cmd))
        subprocess.check_call(cmd)
    else:
        raise RuntimeError('convert.py not found; cannot convert PPTX')


def extract_notes(pptx_path: Path):
    slides_notes = []
    if Presentation is None:
        print('python-pptx not available; skipping notes extraction')
        return slides_notes
    prs = Presentation(str(pptx_path))
    for i, slide in enumerate(prs.slides, start=1):
        notes = ''
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text or ''
        except Exception:
            notes = ''
        slides_notes.append(notes)
    return slides_notes


def main():
    p = argparse.ArgumentParser()
    p.add_argument('--pptx', required=True)
    p.add_argument('--db', default='sqlite:///elearning.db')
    p.add_argument('--title', default=None)
    args = p.parse_args()

    pptx = Path(args.pptx)
    if not pptx.exists():
        print('PPTX not found:', pptx)
        return

    out_dir = Path(__file__).parent / 'static' / 'courses' / pptx.stem
    if out_dir.exists():
        print('Removing existing output dir', out_dir)
        shutil.rmtree(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)

    # convert slides
    convert_pptx(pptx, out_dir)

    # extract notes
    notes = extract_notes(pptx)

    # collect slide images in sorted order
    pngs = sorted(out_dir.glob('*.png'))

    engine = create_engine(args.db)
    Base.metadata.create_all(engine)
    Session = sessionmaker(bind=engine)
    session = Session()

    course_title = args.title or (notes and notes[0]) or pptx.stem
    course = Course(title=course_title)
    session.add(course)
    session.commit()

    for i, pimg in enumerate(pngs, start=1):
        lesson_title = pimg.stem
        lesson_notes = notes[i-1] if (i-1) < len(notes) else ''
        # ensure file path relative to static
        rel = f'courses/{pptx.stem}/{pimg.name}'
        lesson = Lesson(course_id=course.id, title=lesson_title, slide_filename=rel, notes=lesson_notes, index=i)
        session.add(lesson)

    session.commit()
    print('Imported course', course.id, course.title)

if __name__ == '__main__':
    main()
