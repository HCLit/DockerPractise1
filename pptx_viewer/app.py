import os
import json
import time
from flask import Flask, render_template, send_from_directory, url_for
from pathlib import Path
import subprocess

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except Exception:
    Presentation = None
    MSO_SHAPE_TYPE = None

app = Flask(__name__)
BASE_DIR = Path(__file__).parent
SLIDES_DIR = BASE_DIR / "slides"
PPTX_FILE = os.environ.get("PPTX_FILE", str(BASE_DIR.parent / "e_learning_platform_architecture.pptx"))


def extract_notes_and_assets(pptx_path: Path, out_dir: Path):
    """Extract speaker notes and slide assets (images) into out_dir/assets and write notes.json"""
    if Presentation is None:
        print("python-pptx not available; cannot extract notes")
        return

    if not pptx_path.exists():
        print("PPTX not found; skipping notes extraction")
        return

    prs = Presentation(str(pptx_path))
    assets_dir = out_dir / "assets"
    assets_dir.mkdir(parents=True, exist_ok=True)

    slides_data = []

    for i, slide in enumerate(prs.slides, start=1):
        notes_text = ""
        try:
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text or ""
        except Exception:
            notes_text = ""

        assets = []
        n = 0
        for shape in slide.shapes:
            # picture shapes expose .image
            if hasattr(shape, "image"):
                img = shape.image
                ext = img.ext or "png"
                n += 1
                fname = f"slide_{i:03d}_img_{n}.{ext}"
                fpath = assets_dir / fname
                try:
                    fpath.write_bytes(img.blob)
                    assets.append(fname)
                except Exception as e:
                    print("Failed to write asset:", e)

        slides_data.append({"index": i, "notes": notes_text, "assets": assets})

    notes_file = out_dir / "notes.json"
    try:
        notes_file.write_text(json.dumps(slides_data, indent=2), encoding="utf-8")
    except Exception as e:
        print("Failed to write notes.json:", e)


def ensure_notes():
    """Ensure notes.json and assets exist and are up-to-date with the PPTX file."""
    pptx = Path(PPTX_FILE)
    notes_file = SLIDES_DIR / "notes.json"
    if not pptx.exists():
        return

    # If notes.json missing or older than pptx, (re)extract
    if not notes_file.exists() or (pptx.stat().st_mtime > notes_file.stat().st_mtime):
        print("Extracting notes & assets from PPTX...")
        extract_notes_and_assets(pptx, SLIDES_DIR)
        # touch notes_file mtime is done in write, but ensure small sleep to allow mtime differences
        time.sleep(0.1)


def ensure_slides():
    # If slides folder is empty, try to convert
    slides = sorted(SLIDES_DIR.glob("*.png"))
    if slides:
        return

    # Attempt conversion using convert.py
    convert_script = BASE_DIR / "convert.py"
    cmd = ["python", str(convert_script), "--pptx", PPTX_FILE, "--out", str(SLIDES_DIR)]
    print("Converting PPTX -> PNG using:", cmd)
    try:
        subprocess.check_call(cmd)
    except subprocess.CalledProcessError as e:
        print("Conversion failed:", e)

    # After converting slides, extract notes/assets
    try:
        ensure_notes()
    except Exception as e:
        print("Notes extraction failed:", e)


@app.route("/notes")
def notes():
    ensure_slides()
    ensure_notes()

    notes_file = SLIDES_DIR / "notes.json"
    if not notes_file.exists():
        slides = []
    else:
        try:
            slides = json.loads(notes_file.read_text(encoding="utf-8"))
        except Exception:
            slides = []

    return render_template("notes.html", slides=slides)


@app.route("/reveal")
def reveal():
    """Serve a Reveal.js deck generated from slides and notes."""
    ensure_slides()
    ensure_notes()

    notes_file = SLIDES_DIR / "notes.json"
    slides_data = []

    if notes_file.exists():
        try:
            notes = json.loads(notes_file.read_text(encoding="utf-8"))
        except Exception:
            notes = []
    else:
        notes = []

    # Build slides list: prefer notes index mapping; fallback to available PNGs
    for entry in notes:
        idx = entry.get("index")
        img_name = f"slide_{idx:03d}.png"
        img_path = SLIDES_DIR / img_name
        if not img_path.exists():
            # try to find any matching png at position
            candidates = sorted([p.name for p in SLIDES_DIR.glob("*.png")])
            img_name = candidates[idx-1] if (idx-1) < len(candidates) else (candidates[0] if candidates else None)

        slides_data.append({"index": idx, "img": img_name, "notes": entry.get("notes", "")})

    # If no notes present, just use found images
    if not slides_data:
        candidates = sorted([p.name for p in SLIDES_DIR.glob("*.png")])
        for i, name in enumerate(candidates, start=1):
            slides_data.append({"index": i, "img": name, "notes": ""})

    return render_template("reveal.html", slides=slides_data)


@app.route("/")
def index():
    ensure_slides()
    ensure_notes()
    import re

    def sort_key(name):
        m = re.search(r"(\d+)", name)
        return int(m.group(1)) if m else name

    images = sorted([p.name for p in SLIDES_DIR.glob("*.png")], key=sort_key)
    return render_template("index.html", images=images)


@app.route("/slides/<path:filename>")
def slides(filename):
    return send_from_directory(SLIDES_DIR, filename)


if __name__ == "__main__":
    # ensure folder exists
    SLIDES_DIR.mkdir(parents=True, exist_ok=True)
    app.run(host="0.0.0.0", port=int(os.environ.get("PORT", 5000)), debug=True)
